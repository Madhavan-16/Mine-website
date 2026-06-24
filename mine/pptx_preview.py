"""Build an HTML slide preview for .pptx uploads (text + embedded images)."""

from __future__ import annotations

import base64
import logging
from html import escape
from io import BytesIO

logger = logging.getLogger(__name__)

_MAX_SLIDES = 50


def _iter_shapes(shape):
    if getattr(shape, "shapes", None):
        for child in shape.shapes:
            yield from _iter_shapes(child)
    yield shape


def _shape_text_chunks(shape) -> list[str]:
    chunks: list[str] = []
    if getattr(shape, "has_text_frame", False) and shape.text_frame:
        text = (shape.text_frame.text or "").strip()
        if text:
            chunks.append(text)
    if getattr(shape, "has_table", False) and shape.table:
        for row in shape.table.rows:
            cells = [escape((cell.text or "").strip()) for cell in row.cells if (cell.text or "").strip()]
            if cells:
                chunks.append(" · ".join(cells))
    return chunks


def _shape_image_data_uri(shape) -> str | None:
    if not getattr(shape, "image", None):
        return None
    try:
        blob = shape.image.blob
        ext = (shape.image.ext or "png").lower()
    except Exception:
        return None
    if not blob:
        return None
    mime = "image/jpeg" if ext in ("jpg", "jpeg") else f"image/{ext}"
    encoded = base64.b64encode(blob).decode("ascii")
    return f"data:{mime};base64,{encoded}"


def render_pptx_preview_html(
    data: bytes,
    *,
    deck_title: str = "",
    max_slides: int = _MAX_SLIDES,
) -> str | None:
    """Return a full HTML document with slide text and embedded pictures, or None if unreadable."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not available; skipping pptx HTML preview")
        return None

    if not data:
        return None

    try:
        prs = Presentation(BytesIO(data))
    except Exception as e:
        logger.info("pptx preview could not open deck: %s", e)
        return None

    slides = list(prs.slides)
    if not slides:
        return None

    total = len(slides)
    shown = slides[:max_slides]
    truncated = total > max_slides
    doc_title = escape((deck_title or "Presentation").strip() or "Presentation")

    parts: list[str] = [
        "<!DOCTYPE html>",
        '<html lang="en"><head><meta charset="utf-8"/>',
        '<meta name="viewport" content="width=device-width, initial-scale=1"/>',
        f"<title>{doc_title} · preview</title>",
        "<style>",
        "body{margin:0;font-family:system-ui,-apple-system,Segoe UI,Roboto,sans-serif;font-size:13px;background:#0f172a;color:#e2e8f0;}",
        ".wrap{box-sizing:border-box;padding:12px 12px 28px;max-width:960px;margin:0 auto;}",
        "header.page-hd{margin-bottom:14px;}",
        "header.page-hd h1{margin:0 0 6px;font-size:1.1rem;font-weight:600;color:#f8fafc;}",
        ".note{font-size:11px;color:#94a3b8;margin:0;line-height:1.45;}",
        ".warn{font-size:11px;color:#fcd34d;background:#422006;border:1px solid #854d0e;border-radius:6px;padding:8px 10px;margin-bottom:12px;}",
        ".slide{background:linear-gradient(145deg,#1e293b,#0f172a);border:1px solid #334155;border-radius:10px;padding:14px 16px;margin-bottom:12px;box-shadow:0 4px 14px rgba(0,0,0,.25);}",
        ".slide-hd{display:flex;align-items:center;gap:8px;margin:0 0 10px;font-size:12px;font-weight:600;color:#38bdf8;text-transform:uppercase;letter-spacing:.05em;}",
        ".slide-index{background:#0c4a6e;color:#e0f2fe;padding:2px 8px;border-radius:999px;font-size:10px;}",
        ".slide-media{display:flex;flex-wrap:wrap;gap:10px;margin:0 0 10px;}",
        ".slide-media img{max-width:100%;height:auto;border-radius:6px;border:1px solid #334155;background:#020617;}",
        ".slide-body{margin:0;padding-left:18px;color:#cbd5e1;line-height:1.5;}",
        ".slide-body li{margin:4px 0;}",
        ".slide-empty{margin:0;color:#64748b;font-style:italic;}",
        "</style></head><body><div class=\"wrap\">",
        "<header class=\"page-hd\">",
        f"<h1>{doc_title}</h1>",
        "<p class=\"note\">Fallback preview — embedded photos are shown; icon fonts and SmartArt may still be missing. "
        "Use visual slide preview when available for full fidelity.</p>",
        "</header>",
    ]

    if truncated:
        parts.append(
            f'<p class="warn">This deck has <strong>{total}</strong> slides; showing the first <strong>{max_slides}</strong>.</p>'
        )

    for i, slide in enumerate(shown, 1):
        bullets: list[str] = []
        images: list[str] = []
        for shape in slide.shapes:
            for nested in _iter_shapes(shape):
                bullets.extend(_shape_text_chunks(nested))
                uri = _shape_image_data_uri(nested)
                if uri and uri not in images:
                    images.append(uri)

        parts.append('<article class="slide">')
        parts.append(
            f'<h2 class="slide-hd"><span class="slide-index">{i} / {total}</span> Slide {i}</h2>'
        )
        if images:
            parts.append('<div class="slide-media">')
            for uri in images:
                parts.append(f'<img src="{uri}" alt="Slide {i} image" loading="lazy"/>')
            parts.append("</div>")
        if bullets:
            parts.append('<ul class="slide-body">')
            for block in bullets:
                for line in block.splitlines():
                    line = line.strip()
                    if line:
                        parts.append(f"<li>{escape(line)}</li>")
            parts.append("</ul>")
        if not bullets and not images:
            parts.append('<p class="slide-empty">No extractable text or images on this slide.</p>')
        parts.append("</article>")

    parts.append("</div></body></html>")
    return "".join(parts)
