"""Best-effort text extraction from uploaded knowledge files for autofill."""

from __future__ import annotations

import io
import re
from pathlib import Path

# Cap extracted plain text to keep memory and response time bounded (~upload limit order).
_MAX_EXTRACT_CHARS = 350_000


def title_from_filename(filename: str) -> str:
    stem = Path(filename or "document").stem
    stem = re.sub(r"[_\s]+", " ", stem).strip()
    return stem or "Untitled"


def _norm_ws(s: str) -> str:
    return re.sub(r"\s+", " ", (s or "").replace("\x00", "")).strip()


def _clip(s: str, max_len: int) -> str:
    s = _norm_ws(s)
    if len(s) <= max_len:
        return s
    cut = s[:max_len]
    if " " in cut:
        cut = cut.rsplit(" ", 1)[0]
    return cut + "…"


def _gist_summary(blob_raw: str, blob_norm: str, max_sentences: int = 5, max_chars: int = 900) -> str:
    """Short intro: a few sentences/lines, not the full document."""
    text = blob_norm or _norm_ws(blob_raw)
    if not text:
        return ""
    if len(text) <= 280:
        return text[:2000]

    raw_sentences = re.split(r"(?<=[.!?])\s+", text)
    sentences: list[str] = []
    for s in raw_sentences:
        s = s.strip()
        if len(s) < 14:
            continue
        sentences.append(s)
        if len(sentences) >= max_sentences:
            break

    if len(sentences) >= 2:
        joined = "\n".join(sentences[:max_sentences])
        if len(joined) > max_chars:
            joined = _clip(joined, max_chars)
        return joined[:2000]

    lines: list[str] = []
    for ln in blob_raw.splitlines():
        t = _norm_ws(ln)
        t = re.sub(r"^#+\s*", "", t).strip()
        if len(t) < 18:
            continue
        lines.append(t)
        if len(lines) >= max_sentences:
            break
    if len(lines) >= 2:
        joined = "\n".join(_clip(l, 260) for l in lines[:max_sentences])
        return joined[:max_chars][:2000]

    return _clip(text, min(max_chars, 700))


def _clean_autofill_title(title: str) -> str:
    t = (title or "").strip()
    t = re.sub(r"^#+\s*", "", t).strip()
    return t[:500]


def _strip_series_header_for_summary(blob_raw: str, blob_norm: str, title: str) -> tuple[str, str]:
    """Remove leading episode / series / title boilerplate so the summary is substantive content only."""
    title_plain = _clean_autofill_title(title or "")
    t_key = title_plain.lower()

    lines = blob_raw.splitlines()
    i = 0
    while i < len(lines):
        raw = lines[i]
        stripped = raw.strip()
        if not stripped:
            i += 1
            continue
        plain = re.sub(r"^#+\s*", "", stripped).strip()
        low = plain.lower()
        long_line = len(plain) > 420

        if not long_line and re.match(r"^#+\s", stripped):
            i += 1
            continue
        if not long_line and re.search(r"\bdomain term of the week\b", low) and len(plain) < 220:
            i += 1
            continue
        if not long_line and re.match(r"^episode\s*[\d.]+\s*[:.\-–]", low, re.I) and len(plain) < 180:
            i += 1
            continue
        if t_key and len(t_key) > 8 and not long_line:
            if low == t_key or (low.startswith(t_key) and len(plain) <= len(t_key) + 45):
                i += 1
                continue
        break

    raw2 = "\n".join(lines[i:]).strip()
    if not raw2:
        raw2 = blob_raw
    norm2 = _norm_ws(raw2)

    # One-line PDF/PPT dumps: drop prefix through "Domain Term of the Week" (and light punctuation).
    norm2 = re.sub(
        r"(?is)^.{0,320}?\bdomain term of the week\b\s*[.\-–:;•]?\s*",
        "",
        norm2,
        count=1,
    ).strip()

    # Remaining leading episode slug on one line
    norm2 = re.sub(
        r"(?is)^#?\s*episode\s*[\d.]+\s*[:.\-–]\s*.{0,160}?(?=\s+[A-Z]|\s+[a-z]{4,}\s)",
        "",
        norm2,
        count=1,
    ).strip()

    if t_key and len(t_key) > 8:
        lo = norm2.lower()
        if lo.startswith(t_key):
            tail = norm2[len(title_plain) :].lstrip()
            tail = re.sub(r"^[\-–—:.;,•·]+", "", tail)
            if len(tail) > 45:
                norm2 = tail

    if len(norm2) < 42:
        return blob_raw, blob_norm
    return raw2, norm2


def _first_line_title(blob: str, fallback: str, max_len: int = 200) -> str:
    for raw in blob.splitlines():
        line = _norm_ws(raw)
        if len(line) >= 6:
            return (line[:max_len] + "…") if len(line) > max_len else line
    return fallback


def _truncate_blob(blob: str) -> str:
    if len(blob) <= _MAX_EXTRACT_CHARS:
        return blob
    return blob[:_MAX_EXTRACT_CHARS] + "\n…[extraction truncated at size limit]"


def extract_pdf(data: bytes, filename: str) -> tuple[str, str, str]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    meta_title = ""
    md = reader.metadata
    if md:
        raw = getattr(md, "title", None)
        if raw:
            meta_title = _norm_ws(str(raw))
    parts: list[str] = []
    total = 0
    for page in reader.pages:
        if total >= _MAX_EXTRACT_CHARS:
            break
        try:
            t = page.extract_text() or ""
        except Exception:
            t = ""
        if t:
            parts.append(t)
            total += len(t)
    blob = "\n".join(parts)
    blob = _truncate_blob(blob)
    blob_norm = _norm_ws(blob)
    title = meta_title or _first_line_title(blob, title_from_filename(filename))
    title = _clean_autofill_title(title)
    raw_s, norm_s = _strip_series_header_for_summary(blob, blob_norm, title)
    summary = _gist_summary(raw_s, norm_s) if norm_s else ""
    return title, summary, blob


def extract_docx(data: bytes, filename: str) -> tuple[str, str, str]:
    from docx import Document

    doc = Document(io.BytesIO(data))
    title = title_from_filename(filename)
    chunks: list[str] = []

    for p in doc.paragraphs:
        txt = _norm_ws(p.text)
        if not txt:
            continue
        style_name = getattr(getattr(p, "style", None), "name", None) or ""
        if style_name.startswith("Title") and len(txt) <= 500:
            title = txt[:500]
        elif style_name.startswith("Heading 1") and len(txt) <= 500 and title == title_from_filename(filename):
            title = txt[:500]
        chunks.append(txt)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                tx = _norm_ws(cell.text)
                if tx:
                    chunks.append(tx)

    blob = "\n".join(chunks)
    blob = _truncate_blob(blob)
    blob_norm = _norm_ws(blob)

    if title == title_from_filename(filename) and chunks:
        first = chunks[0]
        if 6 <= len(first) <= 300:
            title = first[:500]
    title = _clean_autofill_title(title)

    raw_s, norm_s = _strip_series_header_for_summary(blob, blob_norm, title)
    summary = _gist_summary(raw_s, norm_s) if norm_s else ""
    return title, summary, blob


def _pptx_iter_shapes(shapes):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _pptx_iter_shapes(shape.shapes)
        else:
            yield shape


def _pptx_slide_blocks(data: bytes) -> list[list[str]]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides: list[list[str]] = []
    total = 0
    for slide in prs.slides:
        blocks: list[str] = []
        if total >= _MAX_EXTRACT_CHARS:
            slides.append(blocks)
            continue
        for shape in _pptx_iter_shapes(slide.shapes):
            if hasattr(shape, "text"):
                t = _norm_ws(getattr(shape, "text", "") or "")
                if t:
                    blocks.append(t)
                    total += len(t)
        if slide.has_notes_slide:
            try:
                nf = slide.notes_slide.notes_text_frame
                if nf and nf.text:
                    t = _norm_ws(nf.text)
                    if t:
                        blocks.append(t)
                        total += len(t)
            except Exception:
                pass
        slides.append(blocks)
    return slides


_GENERIC_CASE_STUDY_LABELS = frozenset(
    {
        "case study",
        "case studies",
        "customer profile",
        "benefits",
        "benefit",
        "challenges",
        "challenge",
        "business challenge",
        "solution",
        "our solution",
        "faster",
        "better",
        "cheaper",
        "overview",
        "agenda",
    }
)

_BOILERPLATE_LINE_RE = re.compile(
    r"(?:www\.hexaware\.com|hexaware technologies|all rights reserved|confidential|"
    r"do not distribute|^\s*\d+\s*$|^\s*\|\s*$)",
    re.I,
)

_SIDEBAR_NOISE_RE = re.compile(
    r"(?:\$\s*[\d.]+\s*bn|employees?\s+worldwide|leading\s+stock\s+exchange|"
    r"revenue|headquartered\s+in|customer\s+profile)",
    re.I,
)


def _compact_label(text: str) -> str:
    return re.sub(r"[\s:.\-–—|]+", " ", (text or "").lower()).strip()


def _is_boilerplate_line(line: str) -> bool:
    t = (line or "").strip()
    if not t or len(t) < 3:
        return True
    if _BOILERPLATE_LINE_RE.search(t):
        return True
    if t.count("|") >= 2 and "hexaware" in t.lower():
        return True
    return False


def _is_sidebar_noise(line: str) -> bool:
    return bool(_SIDEBAR_NOISE_RE.search(line or ""))


def _normalize_merged_blob(text: str) -> str:
    """Insert line breaks before section labels, including tight Hexaware one-line exports."""
    t = (text or "").replace("\r\n", "\n").replace("\r", "\n")
    # Slide/page number before section headers: "... distribute. 1 Challenges ..."
    t = re.sub(r"(\.\s*|\s+)\d+\s+(?=(?:business\s+)?challenges?\b)", r"\1\n", t, flags=re.I)
    t = re.sub(r"\s+\d+\s+(?=(?:business\s+)?challenges?\b)", "\n", t, flags=re.I)
    # Column headers often adjacent: "Challenges Solution Faster"
    t = re.sub(r"\b(challenges?)\s+(solution)\b", r"\1\n\2", t, flags=re.I)
    t = re.sub(r"\b(solution)\s+(faster|better|cheaper|benefits?)\b", r"\1\n\2", t, flags=re.I)
    for label in (
        r"(?:business\s+)?challenges?",
        r"solution",
        r"benefits?",
        r"customer\s+profile",
        r"faster",
        r"better",
        r"cheaper",
    ):
        t = re.sub(rf"\s+({label})\s+", rf"\n\1\n", t, flags=re.I)
    t = re.sub(r"^(case\s+stud(?:y|ies))\s+", r"\1\n", t, flags=re.I)
    return t


def _prepare_blob_for_sections(text: str) -> str:
    t = (text or "").replace("\r\n", "\n").replace("\r", "\n")
    if re.search(r"\b(challenges?|solution|benefits?|case\s+study)\b", t, re.I):
        t = _normalize_merged_blob(t)
    return t


def _strip_footer_noise(text: str) -> str:
    t = text or ""
    t = re.split(r"\s+www\.hexaware\.com\b", t, maxsplit=1, flags=re.I)[0]
    t = re.split(r"\|\s*©", t, maxsplit=1, flags=re.I)[0]
    t = re.split(r"\bconfidential\b", t, maxsplit=1, flags=re.I)[0]
    t = re.sub(r"\bdo not distribute\.?\s*", "", t, flags=re.I)
    t = re.sub(r"\ball rights reserved\.?\s*", "", t, flags=re.I)
    t = re.sub(r"\s+\d+\s*$", "", t).strip()
    return t


def _slide_combined_text(blocks: list[str]) -> str:
    return " ".join((b or "").strip() for b in blocks if (b or "").strip())


def _extract_title_from_preamble(preamble: str, filename: str) -> str:
    cleaned = _strip_footer_noise(_strip_case_study_prefix(preamble))
    if not cleaned:
        return ""
    candidates = [_strip_case_study_prefix(x) for x in re.split(r"\n+", _prepare_blob_for_sections(cleaned)) if x.strip()]
    title = _pick_case_study_title(candidates, filename)
    if title and _compact_label(title) not in _GENERIC_CASE_STUDY_LABELS:
        return title
    # Headline often appears before footer noise in the same line.
    head = _strip_footer_noise(cleaned)
    head = _strip_case_study_prefix(head)
    if len(head) >= 20 and _compact_label(head) not in _GENERIC_CASE_STUDY_LABELS:
        return _clean_autofill_title(head[:500])
    return ""


def _find_challenge_bullet_start(raw: str, after: int) -> int:
    m = re.search(
        r"\b(Existing|Data streamed|Lack of|No support|No advanced|The existing|Customer faced|Client faced)\b",
        raw[after:],
        re.I,
    )
    return after + m.start() if m else after


def _find_solution_bullet_start(raw: str, after: int) -> int:
    m = re.search(
        r"\b(Migration|Leveraged|Implemented|Deployed|Introduced|All data other|Azure Blob|"
        r"automated assessment|streaming data analytics)\b",
        raw[after:],
        re.I,
    )
    return after + m.start() if m else len(raw)


def _segment_case_study_blob(text: str, filename: str) -> tuple[str, str, str]:
    """Marker-based split for single-slide decks with column headers in one text run."""
    raw = _prepare_blob_for_sections(text or "")
    if not raw.strip():
        return "", "", ""

    challenge_m = re.search(r"\b(?:business\s+)?challenges?\b", raw, re.I)
    solution_m = re.search(r"\bsolution\b", raw, re.I)
    stop_m = re.search(r"\b(?:faster|better|cheaper|benefits?|customer\s+profile)\b", raw, re.I)

    preamble = raw[: challenge_m.start()] if challenge_m else raw
    title = _extract_title_from_preamble(preamble, filename)

    challenge_body = ""
    solution_body = ""
    stop_at = len(raw)

    if challenge_m:
        # Hexaware-style slides put "Challenges Solution Faster" headers before body text.
        inline_body = ""
        if solution_m and solution_m.start() > challenge_m.end():
            inline_body = raw[challenge_m.end() : solution_m.start()].strip()

        if len(inline_body) >= 40:
            challenge_body = inline_body
            if solution_m:
                solution_body = raw[solution_m.end() : stop_at]
        else:
            ch_start = _find_challenge_bullet_start(raw, challenge_m.end())
            sol_start = _find_solution_bullet_start(raw, ch_start + 20)
            if sol_start > ch_start:
                challenge_body = raw[ch_start:sol_start]
                solution_body = raw[sol_start:stop_at]
            else:
                challenge_body = raw[ch_start:stop_at]
    elif solution_m:
        sol_start = _find_solution_bullet_start(raw, solution_m.end())
        solution_body = raw[sol_start:stop_at]

    # Trim benefits / sidebar tail from solution
    if solution_body:
        ben_in_sol = re.search(r"\b(?:faster|better|cheaper)\b", solution_body, re.I)
        if ben_in_sol and ben_in_sol.start() < 80:
            solution_body = solution_body[ben_in_sol.end() :].strip()

    challenge_body = re.sub(r"^(?:business\s+)?challenges?\s*", "", challenge_body, flags=re.I).strip()
    solution_body = re.sub(r"^solution\s*", "", solution_body, flags=re.I).strip()
    challenge_body = _strip_footer_noise(challenge_body)
    solution_body = _strip_footer_noise(solution_body)

    return (
        title,
        _join_case_study_lines([challenge_body]) if challenge_body else "",
        _join_case_study_lines([solution_body]) if solution_body else "",
    )


def _strip_case_study_prefix(line: str) -> str:
    t = re.sub(r"^case\s+stud(?:y|ies)\s+", "", (line or "").strip(), flags=re.I).strip()
    t = re.split(r"\s+www\.hexaware\.com\b", t, maxsplit=1, flags=re.I)[0].strip()
    t = re.split(r"\s+\|\s*(?:©|all rights reserved|confidential)", t, maxsplit=1, flags=re.I)[0].strip()
    return t


def _explode_text_blocks(blocks: list[str]) -> list[str]:
    """Split large multi-line shapes into lines for section detection."""
    out: list[str] = []
    for block in blocks:
        t = _prepare_blob_for_sections((block or "").strip())
        if not t:
            continue
        needs_split = len(t) > 100 and (
            "\n" in t
            or re.search(r"\b(challenges?|solution|benefits?|faster|better|cheaper)\b", t, re.I)
        )
        if needs_split:
            for piece in re.split(r"\n+", t):
                p = piece.strip()
                if p:
                    out.append(p)
        else:
            out.append(t)
    return out


def _case_study_section_kind(text: str) -> tuple[str | None, str]:
    """Classify a text block as a case-study section header; return (kind, inline body)."""
    raw = (text or "").strip()
    if not raw:
        return None, ""
    compact = _compact_label(raw)

    challenge_labels = {
        "challenge",
        "challenges",
        "business challenge",
        "business challenges",
        "the challenge",
        "the challenges",
        "client challenge",
        "problem statement",
        "the problem",
    }
    solution_labels = {
        "solution",
        "our solution",
        "the solution",
        "proposed solution",
        "approach",
        "our approach",
        "implementation",
    }
    skip_labels = {
        "benefits",
        "benefit",
        "faster",
        "better",
        "cheaper",
        "customer profile",
        "key benefits",
        "outcomes",
        "results",
    }

    if compact in challenge_labels:
        return "challenge", ""
    if compact in solution_labels:
        return "solution", ""
    if compact in skip_labels:
        return "skip", ""

    for pat, kind in (
        (
            re.compile(
                r"^(?:the\s+)?(?:business\s+)?challenges?\s*[:\-–—.]?\s*(.*)$",
                re.I | re.S,
            ),
            "challenge",
        ),
        (
            re.compile(
                r"^(?:client|customer)\s+(?:challenge|situation|context)s?\s*[:\-–—.]?\s*(.*)$",
                re.I | re.S,
            ),
            "challenge",
        ),
        (
            re.compile(
                r"^(?:problem\s+statement|the\s+problem)\s*[:\-–—.]?\s*(.*)$",
                re.I | re.S,
            ),
            "challenge",
        ),
        (
            re.compile(
                r"^(?:our\s+)?(?:proposed\s+)?solution(?:\s+overview)?\s*[:\-–—.]?\s*(.*)$",
                re.I | re.S,
            ),
            "solution",
        ),
        (
            re.compile(r"^(?:our\s+)?approach\s*[:\-–—.]?\s*(.*)$", re.I | re.S),
            "solution",
        ),
        (
            re.compile(r"^(?:key\s+)?benefits?\s*[:\-–—.]?\s*(.*)$", re.I | re.S),
            "skip",
        ),
        (
            re.compile(r"^(?:faster|better|cheaper)\s*[:\-–—.]?\s*(.*)$", re.I | re.S),
            "skip",
        ),
    ):
        m = pat.match(raw)
        if m:
            tail = (m.group(1) or "").strip()
            if tail or len(compact) < 90:
                return kind, tail

    return None, ""


def _pick_case_study_title(candidates: list[str], filename: str) -> str:
    scored: list[tuple[int, str]] = []
    for block in candidates:
        t = _norm_ws(block)
        if not t or _is_boilerplate_line(t):
            continue
        compact = _compact_label(t)
        if compact in _GENERIC_CASE_STUDY_LABELS:
            continue
        if _is_sidebar_noise(t):
            continue
        kind, _ = _case_study_section_kind(t)
        if kind:
            continue
        if len(t) < 12:
            continue
        # Prefer substantive headline-length titles.
        score = len(t)
        if 35 <= len(t) <= 220:
            score += 500
        if re.search(r"\b(for a|transformation|migration|automated|organization|mining)\b", t, re.I):
            score += 200
        if re.search(r"\bmigration from\b|\bleveraged\b|\bloaded in\b|\bdatabricks\b", t, re.I):
            score -= 400
        scored.append((score, t))
    if scored:
        scored.sort(key=lambda x: -x[0])
        return _clean_autofill_title(scored[0][1][:500])
    stem = title_from_filename(filename)
    if _compact_label(stem) in _GENERIC_CASE_STUDY_LABELS:
        return ""
    return _clean_autofill_title(stem)


def _parse_merged_slide_text(text: str) -> tuple[list[str], list[str], list[str]]:
    """Parse one combined shape into title, challenge, and solution line lists."""
    title_lines: list[str] = []
    challenge_lines: list[str] = []
    solution_lines: list[str] = []
    mode = "title"

    prepared = _prepare_blob_for_sections(text or "")
    for line in re.split(r"\n+", prepared):
        line = line.strip()
        if not line or _is_boilerplate_line(line):
            continue

        kind, tail = _case_study_section_kind(line)
        if kind == "challenge":
            mode = "challenge"
            if tail:
                challenge_lines.append(tail)
            continue
        if kind == "solution":
            mode = "solution"
            if tail:
                solution_lines.append(tail)
            continue
        if kind == "skip":
            mode = "skip"
            continue

        compact = _compact_label(line)
        if compact in _GENERIC_CASE_STUDY_LABELS:
            continue
        if _is_sidebar_noise(line):
            continue
        if mode == "skip":
            continue
        if mode == "title":
            stripped = _strip_case_study_prefix(line)
            if stripped and _compact_label(stripped) not in _GENERIC_CASE_STUDY_LABELS:
                title_lines.append(stripped)
        elif mode == "challenge":
            challenge_lines.append(line)
        elif mode == "solution":
            solution_lines.append(line)

    return title_lines, challenge_lines, solution_lines


def _split_bullet_candidates(line: str) -> list[str]:
    """Split run-on PPT paragraphs into separate bullet lines."""
    t = _norm_ws(line)
    if not t:
        return []
    if len(t) < 160:
        return [t]
    parts = re.split(r"(?<=[.!?])\s+(?=[A-Z])", t)
    bullets = [p.strip() for p in parts if len(p.strip()) >= 10]
    return bullets if len(bullets) >= 2 else [t]


def _join_case_study_lines(lines: list[str], max_chars: int = 6000) -> str:
    bullets: list[str] = []
    seen: set[str] = set()
    expanded: list[str] = []
    for line in lines:
        expanded.extend(_split_bullet_candidates(line))
    for line in expanded:
        t = _norm_ws(line)
        if not t or len(t) < 8 or _is_boilerplate_line(t):
            continue
        key = t.lower()
        if key in seen:
            continue
        seen.add(key)
        bullets.append(f"• {t}" if not t.startswith("•") else t)
    if not bullets:
        return ""
    return _clip("\n".join(bullets), max_chars)


def _case_study_title_from_first_slide(blocks: list[str], filename: str) -> str:
    exploded = _explode_text_blocks(blocks)
    return _pick_case_study_title(exploded, filename)


def _join_case_study_parts(parts: list[str], max_chars: int = 6000) -> str:
    cleaned: list[str] = []
    seen: set[str] = set()
    for part in parts:
        t = _norm_ws(part)
        if not t or len(t) < 8:
            continue
        key = t.lower()
        if key in seen:
            continue
        seen.add(key)
        cleaned.append(t)
    if not cleaned:
        return ""
    return _clip("\n\n".join(cleaned), max_chars)


def _slide_body_without_header(blocks: list[str]) -> str:
    if not blocks:
        return ""
    parts: list[str] = []
    for i, block in enumerate(blocks):
        kind, tail = _case_study_section_kind(block)
        if kind and i == 0 and not tail:
            continue
        if kind and i == 0 and tail:
            parts.append(tail)
            continue
        parts.append(block)
    return _join_case_study_parts(parts)


def extract_case_study_pptx(data: bytes, filename: str) -> tuple[str, str, str]:
    """Return title, business challenge, and solution from a case-study deck."""
    slides = _pptx_slide_blocks(data)
    if not slides:
        return title_from_filename(filename), "", ""

    title = ""
    challenge = ""
    solution = ""

    for slide_idx, slide_blocks in enumerate(slides):
        combined = _slide_combined_text(slide_blocks)
        if not combined:
            continue
        t, c, s = _segment_case_study_blob(combined, filename if slide_idx == 0 else "")
        if not title and t:
            title = t
        if c:
            challenge = f"{challenge}\n{c}".strip() if challenge else c
        if s:
            solution = f"{solution}\n{s}".strip() if solution else s

    if not title:
        title = _case_study_title_from_first_slide(slides[0], filename)
    if _compact_label(title) in _GENERIC_CASE_STUDY_LABELS:
        title = _extract_title_from_preamble(_slide_combined_text(slides[0]), filename)

    if not challenge and len(slides) >= 2:
        _, c, _ = _segment_case_study_blob(_slide_combined_text(slides[1]), filename)
        challenge = c
    if not solution and len(slides) >= 3:
        _, _, s = _segment_case_study_blob(_slide_combined_text(slides[2]), filename)
        solution = s

    if challenge and challenge.lower() == title.lower():
        challenge = ""
    if solution and solution.lower() == title.lower():
        solution = ""

    return title[:500], challenge[:6000], solution[:6000]


def extract_pptx(data: bytes, filename: str) -> tuple[str, str, str]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    chunks: list[str] = []
    total = 0

    for slide in prs.slides:
        if total >= _MAX_EXTRACT_CHARS:
            break
        for shape in _pptx_iter_shapes(slide.shapes):
            if hasattr(shape, "text"):
                t = _norm_ws(getattr(shape, "text", "") or "")
                if t:
                    chunks.append(t)
                    total += len(t)
        if slide.has_notes_slide:
            try:
                nf = slide.notes_slide.notes_text_frame
                if nf and nf.text:
                    t = _norm_ws(nf.text)
                    if t:
                        chunks.append(t)
                        total += len(t)
            except Exception:
                pass

    blob = "\n".join(chunks)
    blob = _truncate_blob(blob)
    blob_norm = _norm_ws(blob)

    title = title_from_filename(filename)
    for t in chunks:
        if 6 <= len(t) <= 220:
            title = t[:500]
            break
    title = _clean_autofill_title(title)

    raw_s, norm_s = _strip_series_header_for_summary(blob, blob_norm, title)
    summary = _gist_summary(raw_s, norm_s) if norm_s else ""
    return title, summary, blob


def extract_xlsx(data: bytes, filename: str) -> tuple[str, str, str]:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    try:
        rows: list[str] = []
        total = 0
        max_rows_per_sheet = 8000
        for sheetname in wb.sheetnames:
            if total >= _MAX_EXTRACT_CHARS:
                break
            ws = wb[sheetname]
            rows.append(f"[{sheetname}]")
            for row in ws.iter_rows(max_row=max_rows_per_sheet, values_only=True):
                if total >= _MAX_EXTRACT_CHARS:
                    break
                cells = [_norm_ws(str(c)) for c in row if c is not None and _norm_ws(str(c))]
                if cells:
                    line = " ".join(cells)
                    rows.append(line)
                    total += len(line)
        blob = "\n".join(rows)
        blob = _truncate_blob(blob)
        blob_norm = _norm_ws(blob)

        non_hdr = [r for r in rows if r and not r.startswith("[")]
        title = title_from_filename(filename)
        if non_hdr:
            cand = non_hdr[0]
            if len(cand) <= 200:
                title = cand[:500]
            elif len(non_hdr) > 1:
                cand2 = non_hdr[1]
                if len(cand2) <= 200:
                    title = cand2[:500]

        title = _clean_autofill_title(title[:500])
        raw_s, norm_s = _strip_series_header_for_summary(blob, blob_norm, title)
        summary = _gist_summary(raw_s, norm_s) if norm_s else ""
        return title, summary, blob
    finally:
        wb.close()


def extract_document_text(filename: str, data: bytes) -> tuple[str, str, str]:
    """Return (title, summary, body) for indexing (PDF/DOCX/PPTX/XLSX/TXT/MD)."""
    name = filename or "file"
    ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
    title = title_from_filename(name)
    if ext == "pdf":
        return extract_pdf(data, name)
    if ext == "docx":
        return extract_docx(data, name)
    if ext == "pptx":
        return extract_pptx(data, name)
    if ext == "xlsx":
        return extract_xlsx(data, name)
    if ext in ("txt", "md", "csv", "log"):
        try:
            blob = data.decode("utf-8")
        except UnicodeDecodeError:
            blob = data.decode("latin-1", errors="ignore")
        blob = _truncate_blob(blob)
        blob_norm = _norm_ws(blob)
        title = _clean_autofill_title(_first_line_title(blob, title))
        summary = _gist_summary(blob, blob_norm) if blob_norm else ""
        return title, summary, blob
    return title, "", ""


def suggest_from_upload(filename: str, data: bytes, *, module: str | None = None) -> dict[str, str]:
    """Return autofill fields from an uploaded file (shape depends on knowledge module)."""
    name = filename or "file"
    ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
    title = title_from_filename(name)
    summary = ""
    mod = (module or "").strip()

    if mod == "case_study" and ext == "pptx":
        try:
            cs_title, challenge, solution = extract_case_study_pptx(data, name)
            return {
                "title": _clean_autofill_title((cs_title or title)[:500]),
                "summary": "",
                "body": "",
                "business_challenge": (challenge or "")[:6000],
                "solution": (solution or "")[:6000],
            }
        except Exception:
            pass

    try:
        if ext == "pdf":
            title, summary, _ = extract_pdf(data, name)
        elif ext == "docx":
            title, summary, _ = extract_docx(data, name)
        elif ext == "pptx":
            title, summary, _ = extract_pptx(data, name)
        elif ext == "xlsx":
            title, summary, _ = extract_xlsx(data, name)
        elif ext == "ppt":
            title = title_from_filename(name)
            summary = (
                "Legacy .ppt files cannot be read for auto-fill; use .pptx or type fields manually."
            )
        elif ext in ("png", "jpg", "jpeg"):
            title = title_from_filename(name)
            summary = "Image attachment — add a short description if needed."
        else:
            title = title_from_filename(name)
    except Exception:
        title = title_from_filename(name)
        summary = ""

    result = {
        "title": _clean_autofill_title((title or title_from_filename(name))[:500]),
        "summary": (summary or "")[:2000],
        "body": "",
        "business_challenge": "",
        "solution": "",
    }
    if mod == "case_study" and ext == "pptx":
        slides = _pptx_slide_blocks(data)
        combined = _slide_combined_text(slides[0]) if slides else ""
        if combined:
            seg_title, seg_challenge, seg_solution = _segment_case_study_blob(combined, name)
            if seg_title:
                result["title"] = _clean_autofill_title(seg_title[:500])
            if seg_challenge:
                result["business_challenge"] = seg_challenge[:6000]
            if seg_solution:
                result["solution"] = seg_solution[:6000]
    elif mod == "case_study":
        result["business_challenge"] = (summary or "")[:6000]
    return result
