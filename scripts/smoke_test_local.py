"""Quick local smoke test before deploy."""
import re
import sys
from io import BytesIO

from pptx import Presentation

from mine import create_app


def main() -> int:
    app = create_app()
    client = app.test_client()
    errors: list[str] = []

    for path in ("/", "/login"):
        r = client.get(path)
        if r.status_code not in (200, 302):
            errors.append(f"GET {path} -> {r.status_code}")

    with client.session_transaction() as sess:
        sess["user_id"] = 1

    r = client.get("/content/create?module=case_study")
    if r.status_code != 200:
        errors.append(f"GET /content/create -> {r.status_code}")
    html = r.get_data(as_text=True)
    if "business_challenge" not in html:
        errors.append("case study form fields missing on create page")

    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[5])
    s1.shapes.title.text = "Test Case"
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "Business Challenge"
    s2.shapes.add_textbox(0, 0, 0, 0).text_frame.text = "Operational problem."
    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    s3.shapes.title.text = "Solution"
    s3.shapes.add_textbox(0, 0, 0, 0).text_frame.text = "Analytics platform deployed."
    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)

    m = re.search(r'name="csrf_token" value="([^"]+)"', html)
    if not m:
        errors.append("csrf token missing on create page")
    else:
        r2 = client.post(
            "/content/suggest-fields",
            data={"csrf_token": m.group(1), "module": "case_study", "file": (bio, "case.pptx")},
            content_type="multipart/form-data",
        )
        if r2.status_code != 200:
            errors.append(f"POST /content/suggest-fields -> {r2.status_code}")
        else:
            j = r2.get_json()
            if not j or not j.get("ok"):
                errors.append(f"suggest-fields failed: {j}")
            elif not j.get("business_challenge") or not j.get("solution"):
                errors.append(f"suggest-fields incomplete: {j}")

    r3 = client.get("/files/12/preview-pptx-html")
    if r3.status_code != 200:
        errors.append(f"GET preview-pptx-html -> {r3.status_code}")

    if errors:
        print("SMOKE TEST FAILED:")
        for e in errors:
            print(" -", e)
        return 1
    print("SMOKE TEST OK")
    return 0


if __name__ == "__main__":
    sys.exit(main())
